import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# --- Generate sample test logs ---
np.random.seed(0)
tests = pd.DataFrame({
    'TestID': range(1, 51),
    'Result': np.random.choice(['PASS', 'FAIL'], 50, p=[0.8, 0.2]),
    'Duration_s': np.random.randint(20, 300, 50)
})

tests.to_csv('test_log.csv', index=False)

# --- Read the CSV ---
df = pd.read_csv('test_log.csv')

# --- Summary statistics ---
summary = df['Result'].value_counts(normalize=True) * 100
avg_duration = df['Duration_s'].mean()

print("Test Summary:")
print(summary)
print(f"Average Duration: {avg_duration:.2f} seconds")

# --- Pie chart ---
plt.figure(figsize=(5,5))
plt.pie(summary, labels=summary.index, autopct='%1.1f%%', startangle=90, colors=['#8BC34A','#F44336'])
plt.title('Test Results Distribution')
plt.savefig('test_results_pie.png')
plt.close()

# --- Create PowerPoint report ---
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
slide.shapes.title.text = "Test Automation Summary"

# Add text box for summary
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
tf = txBox.text_frame
tf.text = f"Total Tests: {len(df)}\nPASS: {summary['PASS']:.1f}%\nFAIL: {summary['FAIL']:.1f}%\nAverage Duration: {avg_duration:.2f}s"

# Add pie chart image
slide.shapes.add_picture('test_results_pie.png', Inches(2), Inches(3), width=Inches(6))

# Save PPTX
prs.save('Test_Automation_Report.pptx')
print("PowerPoint report generated: Test_Automation_Report.pptx")
